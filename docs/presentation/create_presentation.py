from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Case Intake System"
    subtitle.text = "Technical Overview"
    
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)

def create_tech_stack_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Tech Stack & Rationale"
    
    content.text = """Frontend:
• React with TypeScript
• Material-UI
• Axios

Backend:
• .NET 7 Web API
• Entity Framework Core
• SQLite

Rationale:
• Modern, maintainable stack
• Strong typing throughout
• Separation of concerns
• Easy to deploy and scale"""

def create_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Architecture Overview"
    
    content.text = """Clean Architecture:
• Core (domain models, interfaces)
• Infrastructure (data access, services)
• API (controllers, endpoints)
• Web (React frontend)

Key Components:
• Case Management System
• File Upload Service
• Role-based Authentication
• API Gateway Pattern"""

def create_database_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Database Structure"
    
    content.text = """Core Entities:
• Cases (id, fullName, email, description, type, status)
• FileAttachments (id, fileName, contentType, filePath)
• Users (id, username, role)

Design Decisions:
• SQLite for development simplicity
• File-based storage for attachments
• Soft delete pattern
• Audit fields (createdAt, updatedAt)"""

def create_assumptions_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Assumptions & Next Steps"
    
    content.text = """Assumptions:
• Single tenant system
• Basic role-based access (admin/user)
• File size limits manageable
• No real-time requirements

Next Steps:
• Implement Microsoft Dynamics integration
• Add CSV export functionality
• Enhance file handling
• Add email notifications"""

def create_production_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Production Improvements"
    
    content.text = """Security:
• JWT authentication
• HTTPS enforcement
• Rate limiting
• Input validation

Performance:
• Caching layer
• File compression
• Database indexing
• API response optimization

Monitoring:
• Logging
• Error tracking
• Performance metrics
• Health checks"""

def create_bonus_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Bonus Features"
    
    content.text = """Login System:
• Role-based access control
• Session management
• Password hashing

Admin Features:
• Case listing with filters
• Bulk operations
• Export to CSV

Integration:
• Microsoft Dynamics payload format
• Webhook support
• API documentation"""

def main():
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    create_title_slide(prs)
    create_tech_stack_slide(prs)
    create_architecture_slide(prs)
    create_database_slide(prs)
    create_assumptions_slide(prs)
    create_production_slide(prs)
    create_bonus_slide(prs)
    
    prs.save('docs/presentation/CaseIntake_Technical_Overview.pptx')

if __name__ == '__main__':
    main() 